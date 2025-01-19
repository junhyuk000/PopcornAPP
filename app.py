from pptx import Presentation
from pptx.util import Pt

def create_movieinfo_ppt(file_name="MovieInfoAPP_Presentation.pptx"):
    ppt = Presentation()

    # Slide 1: Title Slide
    slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    slide.shapes.title.text = "MovieInfoAPP"
    slide.placeholders[1].text = "영화 리뷰와 정보를 한눈에, MovieInfoAPP"

    # Slide 2: Project Overview
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "프로젝트 개요"
    slide.placeholders[1].text = (
        "1. 실시간 영화 API를 활용하여 영화 정보를 시각적으로 제공\n"
        "2. 회원가입 및 로그인 기능\n"
        "3. 리뷰 작성 및 관리 기능\n"
        "4. Kakaomap API를 활용한 영화관 검색\n"
        "5. YouTube API를 활용한 예고편 시청\n"
    )

    # Slide 3: Key Features
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "주요 기능"
    slide.placeholders[1].text = (
        "- 실시간 박스오피스 정보 제공\n"
        "- 리뷰 작성 및 추천 기능\n"
        "- 위치 기반 영화관 검색\n"
        "- 유튜브 예고편 보기\n"
        "- 데이터베이스를 활용한 데이터 관리 및 시각화"
    )

    # Slide 4: Technology Stack
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "사용된 기술"
    slide.placeholders[1].text = (
        "Backend: Python, Flask\n"
        "Frontend: HTML, CSS, JavaScript\n"
        "Database: MySQL\n"
        "Libraries: Pandas, Requests, BeautifulSoup\n"
        "APIs: Kakaomap API, YouTube API"
    )

    # Slide 5: Database Design
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "데이터베이스 설계"
    slide.placeholders[1].text = (
        "1. users 테이블\n"
        "   - idx, id, name, password, user_ip, filename, reg_date\n"
        "2. posts 테이블\n"
        "   - id, userid, username, title, content, rating, spoiler, filename, movie_title, created_at, updated_at, views\n"
        "3. movies_info 테이블\n"
        "   - id, rank, title, genres, director, nations, t_audience, c_audience, t_sales, c_sales, filename, release_date, input_date\n"
        "4. movies 테이블\n"
        "   - id, rank, title, genres, director, nations, t_audience, c_audience, t_sales, c_sales, filename, release_date, input_date"
    )

    # Slide 6: Portfolio Structure
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "프로젝트 구조"
    slide.placeholders[1].text = (
        "real_time_movie_status/\n"
        "   - movie.py: 라우트 연결\n"
        "   - models.py: 데이터베이스 함수 및 API 연동\n"
        "   - templates/: HTML 템플릿 디렉토리\n"
        "       - movie_base.html, movie_movies.html 등\n"
        "       - 카카오 지도 및 유튜브 API 활용 템플릿 포함"
    )

    # Slide 7: Core Functions and APIs
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "주요 함수 및 API"
    slide.placeholders[1].text = (
        "1. movie.py\n"
        "   - movies(): 상영 중 영화 정보\n"
        "   - movie_map(): 카카오 지도 영화관 검색\n"
        "   - movie_youtube(title): 유튜브 예고편 보기\n"
        "2. models.py\n"
        "   - movies_images(): 영화 이미지 저장\n"
        "   - moives_info(): 박스오피스 API 데이터 수집\n"
        "   - insert_data(): 데이터베이스 저장"
    )

    # Slide 8: Results and Future Work
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "프로젝트 결과"
    slide.placeholders[1].text = (
        "- 사용자 친화적인 영화 리뷰 및 정보 제공\n"
        "- 실시간 API 데이터 연동 경험\n"
        "- 데이터 시각화 및 사용자 경험 향상\n"
        "- 향후 개선점:\n"
        "  1. 리뷰 추천 알고리즘 고도화\n"
        "  2. 영화 랭킹 데이터 시각화 고도화"
    )

    # Slide 9: Thank You
    slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    slide.shapes.title.text = "감사합니다"
    slide.placeholders[1].text = "질문이 있으시면 말씀해 주세요!"

    # Save the presentation
    ppt.save(file_name)
    print(f"Presentation '{file_name}' created successfully.")

create_movieinfo_ppt()
